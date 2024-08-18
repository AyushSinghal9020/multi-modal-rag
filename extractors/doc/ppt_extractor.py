from pptx import Presentation

def pptx_to_text(path) : 

    pr = Presentation(path)
    text = ''

    for slide in pr.slides : 

        for shape in slide.shapes : 
            
            if shape.has_text_frame : 
                
                for paragraph in shape.text_frame.paragraphs : 
                    
                    for run in paragraph.runs : text += run.text

    return text